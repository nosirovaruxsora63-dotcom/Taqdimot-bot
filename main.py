import asyncio
import logging
import os
from aiogram import Bot, Dispatcher, F, html
from aiogram.client.default import DefaultBotProperties
from aiogram.enums import ParseMode
from aiogram.filters import CommandStart
from aiogram.types import BufferedInputFile, Message
from dotenv import load_dotenv
from groq import Groq
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

load_dotenv()

TOKEN = "AQ.Ab8RN6Iu6riqs8tEMVX_qo83q8CcYRaNLrCNliWfF9tz4fD3mg"
GROQ_API_KEY = "gsk_... (Groq saytidan olgan kalitingizni yozing)"

bot = Bot(
    token=TOKEN, default=DefaultBotProperties(parse_mode=ParseMode.HTML)
)
dp = Dispatcher()
groq_client = Groq(api_key=GROQ_API_KEY)


def generate_presentation_content(topic: str) -> str:
  prompt = f"""
    Siz professional taqdimot (slayd) dizayneri va mutaxassisisiz. 
    Quyidagi mavzu bo'yicha ANIQ 15 ta slayd uchun mukammal tuzilma tayyorlab bering.
    Har bir slayd uchun quyidagi formatdan aniq foydalaning:
    
    SLAYD N
    Sarlavha: [Slayd sarlavhasi]
    Matn:
    - 1-reja (asosiy fikr yoki qadam)
    - 2-reja (tavsif yoki tahlil)
    - 3-reja (xulosa yoki qo'shimcha ma'lumot)

    Mavzu: {topic}
    """

  chat_completion = groq_client.chat.completions.create(
      messages=[{
          "role": "user",
          "content": prompt,
      }],
      model="llama-3.3-70b-versatile",
      temperature=0.7,
  )
  return chat_completion.choices[0].message.content


def create_pptx_file(content: str, filename: str = "presentation.pptx"):
  prs = Presentation()
  # Slayd o'lchamini kengaytirish (16:9 format zamonaviyroq ko'rinadi)
  prs.slide_width = Inches(13.333)
  prs.slide_height = Inches(7.5)

  slides_data = content.split("SLAYD")

  for block in slides_data:
    if not block.strip():
      continue

    lines = [line.strip() for line in block.split("\n") if line.strip()]
    title = "Slayd"
    bullet_points = []

    for line in lines:
      if line.lower().startswith("sarlavha:"):
        title = line.split(":", 1)[1].strip()
      elif line.startswith("-") or line.startswith("*"):
        bullet_points.append(line.lstrip("-* ").strip())

    # Bo'sh slayd qo'shish (Blank layout)
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 1. Orqa fon rangini och tusda qilish
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 247, 250)  # Och kulrang-moviy fon

    # 2. Yuqori qismga chiroyli dizayn paneli (Banner) qo'shish
    header_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.8), Inches(11.733), Inches(1.1)
    )
    header_box.fill.solid()
    header_box.fill.fore_color.rgb = RGBColor(26, 54, 93)  # To'q ko'k rang
    header_box.line.color.rgb = RGBColor(26, 54, 93)

    # 3. Sarlavha matnini qo'shish
    title_box = slide.shapes.add_textbox(
        Inches(1.0), Inches(0.95), Inches(11.333), Inches(0.8)
    )
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = title
    p_title.font.size = Pt(28)
    p_title.font.bold = True
    p_title.font.color.rgb = RGBColor(255, 255, 255)  # Oq rangli sarlavha

    # 4. Matnlar (Rejalar) uchun asosiy quti
    content_box = slide.shapes.add_textbox(
        Inches(1.0), Inches(2.3), Inches(11.333), Inches(4.5)
    )
    tf_content = content_box.text_frame
    tf_content.word_wrap = True

    # Rejalarni joylashtirish (kamida 3 ta reja bo'lishi ta'minlanadi)
    if not bullet_points:
      bullet_points = [
          "Mavzu bo'yicha asosiy tushunchalar",
          "Asosiy jarayonlar va tahlillar",
          "Amaliy ahamiyati va xulosalar",
      ]

    for i, point in enumerate(bullet_points):
      if i == 0:
        p = tf_content.paragraphs[0]
      else:
        p = tf_content.add_paragraph()

      p.text = f"🔹  {point}"
      p.level = 0
      p.font.size = Pt(20)
      p.font.color.rgb = RGBColor(45, 55, 72)  # To'q kulrang matn
      p.space_after = Pt(14)  # Qatorlar orasidagi masofa

  prs.save(filename)
  return filename


@dp.message(CommandStart())
async def command_start_handler(message: Message):
  await message.answer(
      f"Salom, {html.bold(message.from_user.full_name)}! 🚀\nMen"
      " tayyorman. Menga istalgan mavzuni yuboring, men Groq intellekti"
      " yordamida **15 ta slayddan iborat, chiroyli dizaynli** taqdimot"
      " tayyorlab beraman."
  )


@dp.message(F.text)
async def handle_topic(message: Message):
  topic = message.text
  processing_msg = await message.answer(
      "✨ Groq intellekti 15 ta slayd mazmunini va dizaynini"
      " shakllantirmoqda, iltimos biroz kuting..."
  )

  try:
    # Groq orqali matn generatsiya qilish
    ai_content = await asyncio.to_thread(
        generate_presentation_content, topic
    )

    # PPTX fayl yaratish
    file_path = f"presentation_{message.from_user.id}.pptx"
    await asyncio.to_thread(create_pptx_file, ai_content, file_path)

    # Faylni yuborish
    with open(file_path, "rb") as f:
      file_bytes = f.read()

    document = BufferedInputFile(
        file_bytes, filename=f"{topic[:20]}_taqdimot.pptx"
    )
    await message.answer_document(
        document=document,
        caption=(
            f"<b>{html.escape(topic)}</b> mavzusida 15 ta slayddan iborat tayyor"
            " taqdimot! 🎯"
        ),
    )

    # Vaqtinchalik faylni o'chirish
    if os.path.exists(file_path):
      os.remove(file_path)

    await bot.delete_message(
        chat_id=message.chat.id, message_id=processing_msg.message_id
    )

  except Exception as e:
    logging.error(f"Xatolik: {e}")
    await message.answer(
        "Kechirasiz, taqdimotni yaratishda xatolik yuz berdi. Iltimos, boshqa"
        " mavzu yuboring yoki qaytadan urinib ko'ring."
    )


async def main():
  logging.basicConfig(level=logging.INFO)
  await dp.start_polling(bot)


if __name__ == "__main__":
  asyncio.run(main())
